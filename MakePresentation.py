from pptx import Presentation

# Create a PowerPoint presentation object
prs = Presentation()

# Slide 1: Title Slide
slide_1 = prs.slides.add_slide(prs.slide_layouts[0])
title_1 = slide_1.shapes.title
subtitle_1 = slide_1.placeholders[1]
title_1.text = "How to Establish Key Focus of Analysis for Data Analysis and Visualization Projects"
subtitle_1.text = "Key Steps, Techniques, and Real-Life Examples"

# Slide 2: Definition Slide
slide_2 = prs.slides.add_slide(prs.slide_layouts[1])
title_2 = slide_2.shapes.title
title_2.text = "Definition: Key Focus of Analysis"
content_2 = slide_2.shapes.placeholders[1].text_frame
content_2.text = "Key focus of analysis is the conduct of focused analysis using data analysis to derive business insights."

# Slide 3: Step 1 - Define the Problem or Opportunity
slide_3 = prs.slides.add_slide(prs.slide_layouts[1])
title_3 = slide_3.shapes.title
title_3.text = "Step 1: Define the Problem or Opportunity"
content_3 = slide_3.shapes.placeholders[1].text_frame
content_3.text = "1. Identify a specific business problem or opportunity.\n" \
                 "2. Clarify the objectives and key performance indicators (KPIs).\n" \
                 "3. Develop a hypothesis or research question."

# Slide 4: Step 2 - Data Collection and Preparation
slide_4 = prs.slides.add_slide(prs.slide_layouts[1])
title_4 = slide_4.shapes.title
title_4.text = "Step 2: Data Collection and Preparation"
content_4 = slide_4.shapes.placeholders[1].text_frame
content_4.text = "1. Gather relevant data from various sources (e.g., databases, spreadsheets, surveys).\n" \
                 "2. Clean and preprocess data (handle missing values, outliers, and data formatting).\n" \
                 "3. Transform data into a suitable format for analysis."

# Slide 5: Step 3 - Exploratory Data Analysis (EDA)
slide_5 = prs.slides.add_slide(prs.slide_layouts[1])
title_5 = slide_5.shapes.title
title_5.text = "Step 3: Exploratory Data Analysis (EDA)"
content_5 = slide_5.shapes.placeholders[1].text_frame
content_5.text = "1. Visualize data using plots, charts, and heatmaps.\n" \
                 "2. Summarize data using statistical measures (mean, median, standard deviation).\n" \
                 "3. Identify patterns, trends, and correlations."

# Slide 6: Step 4 - Focus Analysis
slide_6 = prs.slides.add_slide(prs.slide_layouts[1])
title_6 = slide_6.shapes.title
title_6.text = "Step 4: Focus Analysis"
content_6 = slide_6.shapes.placeholders[1].text_frame
content_6.text = "1. Identify key variables and relationships.\n" \
                 "2. Apply statistical techniques (regression, clustering, decision trees).\n" \
                 "3. Use data visualization to communicate findings."

# Slide 7: Common Focus Analysis Techniques
slide_7 = prs.slides.add_slide(prs.slide_layouts[1])
title_7 = slide_7.shapes.title
title_7.text = "Common Focus Analysis Techniques"
content_7 = slide_7.shapes.placeholders[1].text_frame
content_7.text = "1. Correlation Analysis: Identify relationships between variables.\n" \
                 "2. Regression Analysis: Model relationships between variables.\n" \
                 "3. Cluster Analysis: Group similar data points.\n" \
                 "4. Decision Tree Analysis: Identify decision-making patterns.\n" \
                 "5. Segmentation Analysis: Divide customers into segments."

# Slide 8: Step 5 - Insight Generation
slide_8 = prs.slides.add_slide(prs.slide_layouts[1])
title_8 = slide_8.shapes.title
title_8.text = "Step 5: Insight Generation"
content_8 = slide_8.shapes.placeholders[1].text_frame
content_8.text = "1. Interpret results in the context of the business problem.\n" \
                 "2. Identify key drivers, opportunities, and challenges.\n" \
                 "3. Develop actionable recommendations."

# Slide 9: Step 6 - Communication & Storytelling
slide_9 = prs.slides.add_slide(prs.slide_layouts[1])
title_9 = slide_9.shapes.title
title_9.text = "Step 6: Communication and Storytelling"
content_9 = slide_9.shapes.placeholders[1].text_frame
content_9.text = "1. Create clear, concise reports and visualizations.\n" \
                 "2. Use narrative techniques to convey insights.\n" \
                 "3. Present findings to stakeholders."

# Slide 10: Tools and Techniques
slide_10 = prs.slides.add_slide(prs.slide_layouts[1])
title_10 = slide_10.shapes.title
title_10.text = "Tools and Techniques"
content_10 = slide_10.shapes.placeholders[1].text_frame
content_10.text = "1. Data visualisation tools (Tableau, Power BI, D3.js, QlikSense).\n" \
                  "2. Statistical software (R, Python, SPSS).\n" \
                  "3. Machine learning libraries (scikit-learn, TensorFlow).\n" \
                  "4. Data mining tools (SQL, Excel)."

# Slide 11: Real-Life Business Problems
slide_11 = prs.slides.add_slide(prs.slide_layouts[1])
title_11 = slide_11.shapes.title
title_11.text = "Real-Life Business Problems"
content_11 = slide_11.shapes.placeholders[1].text_frame
content_11.text = "1. Customer churn prediction.\n" \
                  "2. Sales & Inventory forecasting.\n" \
                  "3. Market segmentation.\n" \
                  "4. Supply chain optimization.\n" \
                  "5. Pricing strategy development."

# Slide 12: Example - Customer Retention
slide_12 = prs.slides.add_slide(prs.slide_layouts[1])
title_12 = slide_12.shapes.title
title_12.text = "Example: Customer Retention Problem"
content_12 = slide_12.shapes.placeholders[1].text_frame
content_12.text = "Problem: A retail company wants to improve customer retention.\n\n" \
                  "Focus Analysis:\n" \
                  "1. Collect customer data (demographics, purchase history).\n" \
                  "2. Apply clustering analysis to identify customer segments.\n" \
                  "3. Use regression analysis to model retention predictors.\n" \
                  "4. Visualize findings using heatmaps and scatter plots.\n\n" \
                  "Insights:\n" \
                  "1. Identified high-value customer segments.\n" \
                  "2. Discovered key retention drivers (purchase frequency, customer service).\n" \
                  "3. Recommended targeted marketing campaigns."

# Save the presentation to a file
pptx_file = "/mnt/data/Key_Focus_of_Analysis_Data_Projects.pptx"
prs.save(pptx_file)

pptx_file